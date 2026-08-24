from __future__ import annotations

import hashlib
from pathlib import Path

from pptx import Presentation

from .contracts import EvidenceAtom, FactObservation
from .text_extraction import extract_facts_from_text
from .xlsx_parser import load_rules, sha256_file


def parse_pptx(path: Path, root: Path, rules_path: Path) -> tuple[list[EvidenceAtom], list[FactObservation]]:
    rules = load_rules(rules_path)
    document_hash = sha256_file(path)
    relative_path = path.relative_to(root).as_posix()
    presentation = Presentation(path)
    evidence: list[EvidenceAtom] = []
    observations: list[FactObservation] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            shape_id = getattr(shape, "shape_id", 0)
            base_locator = f"slide:{slide_index}/shape:{shape_id}"
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if text:
                    _append_block(text, base_locator, document_hash, relative_path, rules, evidence, observations)
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for column_index, cell in enumerate(row.cells, start=1):
                        text = cell.text.strip()
                        if not text:
                            continue
                        label = row.cells[column_index - 2].text.strip() if column_index > 1 else ""
                        block = f"{label}: {text}" if label else text
                        locator = f"{base_locator}/table:1/row:{row_index}/col:{column_index}"
                        _append_block(block, locator, document_hash, relative_path, rules, evidence, observations, raw_text=text)

    return evidence, observations


def _append_block(
    text: str,
    locator: str,
    document_hash: str,
    relative_path: str,
    rules: dict,
    evidence: list[EvidenceAtom],
    observations: list[FactObservation],
    raw_text: str | None = None,
) -> None:
    evidence_id = hashlib.sha256(f"{document_hash}|{locator}".encode("utf-8")).hexdigest()[:20]
    evidence.append(
        EvidenceAtom(
            evidence_id=evidence_id,
            document_hash=document_hash,
            relative_path=relative_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            locator=locator,
            raw_text=raw_text or text,
            extractor="python-pptx",
        )
    )
    observations.extend(
        extract_facts_from_text(
            text=text,
            base_locator=locator,
            document_hash=document_hash,
            relative_path=relative_path,
            evidence_id=evidence_id,
            rules=rules,
        )
    )

