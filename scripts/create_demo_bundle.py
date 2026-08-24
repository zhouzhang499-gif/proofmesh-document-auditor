from __future__ import annotations

import json
from pathlib import Path

from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
DEMO = ROOT / "examples" / "demo_bundle"


def write_workbook(path: Path, sheet_name: str, rows: list[list[object]]) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    for row in rows:
        sheet.append(row)
    sheet.column_dimensions["A"].width = 22
    sheet.column_dimensions["B"].width = 18
    workbook.save(path)


def main() -> None:
    DEMO.mkdir(parents=True, exist_ok=True)
    write_workbook(
        DEMO / "预算.xlsx",
        "预算表",
        [["指标", "数值"], ["项目总预算", "1260万元"], ["2026Q1同比增长率", "18%"]],
    )
    write_workbook(
        DEMO / "管理层汇报数据.xlsx",
        "汇报数据",
        [["指标", "数值"], ["预算金额", "126万元"], ["2026Q1环比增长率", "18%"]],
    )
    document = Document()
    document.add_heading("项目方案", level=1)
    document.add_paragraph("项目预算为1260万元。2026Q1同比增长率为18%。")
    document.add_paragraph("项目实施费用为300万元。")
    document.save(DEMO / "方案.docx")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = "项目预算为126万元\n2026Q1环比增长率为18%\n项目执行费用为300万元"
    presentation.save(DEMO / "管理层汇报.pptx")

    _write_minimal_pdf(DEMO / "contract-summary.pdf", "Budget: 12600000 CNY")
    ground_truth = {
        "expected_issue_count": 2,
        "expected_review_candidate_count": 1,
        "issues": [
            {"rule_id": "money_value_consistency", "locators": ["预算表!B2", "汇报数据!B2"]},
            {"rule_id": "growth_basis_consistency", "locators": ["预算表!B3", "汇报数据!B3"]},
        ],
        "review_candidates": [
            {"left_label": "项目实施费用", "right_label": "项目执行费用", "status": "needs_review"}
        ],
    }
    (DEMO / "ground_truth.json").write_text(json.dumps(ground_truth, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(DEMO)


def _write_minimal_pdf(path: Path, text: str) -> None:
    safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({safe}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode("ascii") + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii"))
        output.extend(body)
        output.extend(b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    path.write_bytes(bytes(output))


if __name__ == "__main__":
    main()
